import os
import sys
import argparse
from pathlib import Path

if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding='utf-8')
    except AttributeError:
        pass

# Tentativa de importação das bibliotecas de extração
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from abp_tutor.db_app import _get_client
from abp_tutor.logging_setup import logger
from dotenv import load_dotenv

load_dotenv()


def extract_text_from_pdf(path: str) -> str:
    if not fitz:
        raise ImportError("Biblioteca 'pymupdf' nao encontrada. Instale com 'pip install pymupdf'")
    
    text = ""
    with fitz.open(path) as doc:
        for page in doc:
            text += page.get_text() + "\n"
    return text


def extract_text_from_pptx(path: str) -> str:
    if not Presentation:
        raise ImportError("Biblioteca 'python-pptx' nao encontrada. Instale com 'pip install python-pptx'")
    
    prs = Presentation(path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)


def extract_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    
    if ext in [".txt", ".md"]:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    elif ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Extensao de arquivo nao suportada: {ext}")


def upload_material(topic: str, file_path: str):
    if not os.path.exists(file_path):
        print(f"Erro: Arquivo nao encontrado: {file_path}")
        return

    source_file = Path(file_path).name

    print(f"Lendo e extraindo texto de: {file_path}...")
    try:
        content = extract_text(file_path)
        content = content.strip()
        
        if not content:
            print("Erro: O arquivo parece estar vazio ou nao possui texto extraivel.")
            return

        print(f"OK: Texto extraido ({len(content)} caracteres). Enviando para o Supabase...")
        
        client = _get_client()
        data = {
            "macro_topic": topic,
            "source_file": source_file,
            "content": content,
            "updated_at": "now()"
        }
        
        try:
            resp = client.table("tutor_materials").upsert(
                data, on_conflict="macro_topic,source_file"
            ).execute()
            print(f"Sucesso! Material '{source_file}' para '{topic}' salvo no banco.")
        except Exception as e:
            print(f"Erro ao salvar no banco: {e}")

    except Exception as e:
        print(f"Falha na extracao/upload: {e}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Upload de material de estudo para o Tutor ABP")
    parser.add_argument("topic", help="O nome exato do Macro Tema (ex: 'Esquizofrenia')")
    parser.add_argument("file", help="Caminho para o arquivo (PDF, PPTX, TXT ou MD)")
    
    args = parser.parse_args()
    upload_material(args.topic, args.file)
